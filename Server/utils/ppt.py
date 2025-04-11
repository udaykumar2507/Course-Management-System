from pptx import Presentation
from pptx.util import Inches

def add_slide(prs, title, content, bullet_points=True):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    if bullet_points:
        for point in content:
            p = content_placeholder.text_frame.add_paragraph()
            p.text = point
    else:
        content_placeholder.text = content

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "AI-Powered Real-Time Fraud Detection System"
slide.placeholders[1].text = "Protecting Digital Transactions with AI & Blockchain\n\nUday Kumar Maurya\n[Hackathon Name]"

# Slide 2: Problem Statement
add_slide(prs, "Problem Statement", [
    "Fraud in digital transactions is growing rapidly",
    "Traditional fraud detection is rule-based & reactive",
    "Cybercriminals continuously bypass security",
    "Businesses lose billions due to fraudulent transactions",
    "🔴 Need: A real-time AI-driven fraud detection system!"
])

# Slide 3: Our Innovative Solution
add_slide(prs, "Our Innovative Solution", [
    "🟢 AI-Powered Real-Time Fraud Detection",
    "Uses Machine Learning & Anomaly Detection",
    "Monitors transactions in real-time",
    "Assigns a Fraud Risk Score (0-100) to each transaction",
    "Blockchain logs fraud cases securely"
])

# Slide 4: Key Features
add_slide(prs, "Key Features", [
    "✅ AI-Powered Early Warning System",
    "✅ Smart Authentication (Behavioral Biometrics)",
    "✅ Real-Time Transaction Monitoring",
    "✅ Fraud Score System (0-100 Scale)",
    "✅ Blockchain for Secure Fraud Tracking"
])

# Slide 5: How It Works (Step-by-Step)
add_slide(prs, "How It Works?", [
    "1️⃣ Data Collection - Monitors transactions",
    "2️⃣ AI & Anomaly Detection - Flags unusual behavior",
    "3️⃣ Fraud Score System - Scores transactions (0-100)",
    "4️⃣ Real-Time Alerts - Notifies user & blocks fraud",
    "5️⃣ Blockchain for Fraud Tracking - Stores fraud cases securely"
])

# Slide 6: Why is This Better?
add_slide(prs, "Why is This Better?", [
    "🔴 Traditional Systems: Rule-based, static, easy to bypass",
    "🟢 Our AI System: Real-time fraud detection, self-learning AI",
    "✅ Prevents fraud before it happens",
    "✅ Adapts to new fraud patterns",
    "✅ Ensures accurate detection with Fraud Score"
])

# Slide 7: Real-World Impact
add_slide(prs, "Real-World Impact", [
    "💰 For Banks & Fintech - Reduces financial fraud",
    "🛍️ For E-commerce - Stops payment fraud",
    "👥 For Consumers - Protects accounts from scams",
    "🚀 Makes digital transactions safer for everyone!"
])

# Slide 8: Future Scope
add_slide(prs, "Future Scope", [
    "🔹 Advanced AI Models - Improve fraud detection accuracy",
    "🔹 Cross-Industry Collaboration - Shared fraud intelligence",
    "🔹 Expansion to IoT & Blockchain - Secure smart devices"
])

# Slide 9: Conclusion
add_slide(prs, "Conclusion & Call to Action", [
    "Fraud in digital transactions is a growing problem",
    "Our AI-powered fraud detection system provides real-time security",
    "It protects consumers & businesses while being easy to integrate",
    "🚀 Join us in making the digital world safer!"
])

# Slide 10: Q&A
add_slide(prs, "Thank You!", "💬 Any Questions?", bullet_points=False)

# Save the presentation
prs.save("AI_Fraud_Detection.pptx")

print("PowerPoint presentation 'AI_Fraud_Detection.pptx' created successfully!")
